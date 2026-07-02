from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        "title": "Football Match Prediction for FIFA World Cup 2026",
        "content": [
            "Machine learning dashboard for match and tournament forecasts",
            "XGBoost model + Monte Carlo simulation",
            "Streamlit interface for live predictions",
        ],
    },
    {
        "title": "Project scope",
        "content": [
            "Predicts single match outcomes: home win / draw / away win",
            "Simulates the rest of the World Cup using those probabilities",
            "Displays predictions in a Streamlit dashboard",
        ],
    },
    {
        "title": "Architecture and stack",
        "content": [
            "Python",
            "XGBoost for match prediction",
            "Streamlit for dashboard",
            "CSV / SQLite data storage",
            "Monte Carlo simulation for tournament probability aggregation",
        ],
    },
    {
        "title": "Model inputs and features",
        "content": [
            "Home and away Elo ratings",
            "Elo rating difference",
            "Recent form for each team (last 5 matches): wins, draws, losses, goals for/against, goal diff",
            "Derived feature differences",
            "Neutral venue flag",
        ],
    },
    {
        "title": "Dashboard usage",
        "content": [
            "Run `streamlit run app/streamlit_app.py`",
            "Choose an upcoming match from the dropdown",
            "Read the probability table and bar chart",
            "Optionally run the Monte Carlo tournament simulation",
        ],
    },
    {
        "title": "Single match prediction",
        "content": [
            "Home Win (H): probability first team wins",
            "Draw (D): probability of a tied match",
            "Away Win (A): probability second team wins",
            "Example: Curacao vs Ivory Coast → Home 9.42%, Draw 17.09%, Away 73.49%",
        ],
    },
    {
        "title": "Monte Carlo Tournament Simulation",
        "content": [
            "Uses match probabilities from XGBoost",
            "Simulates the remaining tournament many times",
            "Randomly chooses match outcomes based on predicted odds",
            "Aggregates team performance across simulations",
        ],
    },
    {
        "title": "Simulation slider meaning",
        "content": [
            "Higher number = more stable results",
            "Lower number = faster, more variable results",
            "Does not change match probabilities",
            "Just improves the accuracy of tournament probability estimates",
        ],
    },
    {
        "title": "Output metrics explained",
        "content": [
            "win_pct: percent of simulations the team won the tournament",
            "final_pct: percent team reached the final",
            "semi_pct: percent team reached the semifinals",
            "quarter_pct: percent team reached the quarterfinals",
            "r16_pct: percent team reached the Round of 16",
            "r32_pct: percent team reached the Round of 32",
            "wins: raw count of tournament wins in simulations",
        ],
    },
    {
        "title": "Summary and next steps",
        "content": [
            "XGBoost predicts individual match outcomes",
            "Monte Carlo converts match probabilities into tournament forecasts",
            "Dashboard explores match and tournament predictions",
            "Future work: bookie odds, injuries, head-to-head history, live updates",
        ],
    },
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[5]
for i, slide_data in enumerate(slides):
    slide = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_shape.text = slide_data["title"]
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    body = slide.shapes.add_textbox(left, top, width, height)
    tf = body.text_frame
    tf.margin_bottom = Inches(0.05)
    for j, paragraph in enumerate(slide_data["content"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = paragraph
        p.level = 0
        p.font.size = Pt(24)

prs.save("presentation_deck.pptx")
print("presentation_deck.pptx created")
